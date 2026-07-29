from __future__ import annotations

import csv
import io
import logging
import re
import zipfile
import tarfile
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class FileSignals:
    text: str = ""
    metadata: dict = field(default_factory=dict)
    file_type: str = ""
    extension: str = ""
    confidence_hint: float = 0.5
    signals_present: list[str] = field(default_factory=list)


_INSTALLER_EXTENSIONS = frozenset({
    ".exe", ".msi", ".dmg", ".deb", ".rpm", ".appimage",
    ".pkg", ".run", ".sh", ".bin",
})

_SYSTEM_EXTENSIONS = frozenset({
    ".dll", ".sys", ".so", ".dylib", ".kext", ".drv",
    ".o", ".obj", ".lib", ".a",
})

_TEMP_EXTENSIONS = frozenset({
    ".tmp", ".temp", ".swp", ".bak", "~", ".log", ".cache",
})

_CODE_EXTENSIONS = frozenset({
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".rb", ".go", ".rs", ".swift", ".kt", ".scala",
    ".html", ".css", ".scss", ".less", ".xml", ".json", ".yaml",
    ".yml", ".toml", ".ini", ".cfg", ".conf", ".sh", ".bash",
    ".zsh", ".ps1", ".bat", ".cmd", ".sql", ".r", ".m", ".mm",
})

_TEXT_EXTENSIONS = frozenset({
    ".txt", ".md", ".rst", ".tex", ".org",
})

_IMAGE_EXTENSIONS = frozenset({
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif",
    ".webp", ".ico", ".svg", ".heic", ".heif",
})

_AUDIO_EXTENSIONS = frozenset({
    ".mp3", ".wav", ".flac", ".aac", ".ogg", ".wma", ".m4a",
    ".opus", ".aiff", ".wv", ".ape",
})

_VIDEO_EXTENSIONS = frozenset({
    ".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm",
    ".m4v", ".3gp", ".ogv", ".ts",
})

_SPREADSHEET_EXTENSIONS = frozenset({".xlsx", ".xls", ".xlsm", ".xlsb", ".ods"})

_ARCHIVE_EXTENSIONS = frozenset({".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz", ".zst"})


def _get_category(ext: str) -> str:
    if ext in _INSTALLER_EXTENSIONS:
        return "installer"
    if ext in _SYSTEM_EXTENSIONS:
        return "system"
    if ext in _TEMP_EXTENSIONS:
        return "temp"
    if ext in _CODE_EXTENSIONS:
        return "code"
    if ext in _TEXT_EXTENSIONS:
        return "document"
    if ext in _IMAGE_EXTENSIONS:
        return "image"
    if ext in _AUDIO_EXTENSIONS:
        return "audio"
    if ext in _VIDEO_EXTENSIONS:
        return "video"
    if ext in _SPREADSHEET_EXTENSIONS or ext == ".csv":
        return "spreadsheet"
    if ext in _ARCHIVE_EXTENSIONS:
        return "archive"
    if ext == ".pdf":
        return "document"
    if ext in {".docx", ".doc", ".rtf", ".odt", ".pptx", ".ppt"}:
        return "document"
    return "unknown"


def extract_all(file_path: Path) -> FileSignals:
    ext = file_path.suffix.lower()
    category = _get_category(ext)

    if category in ("installer", "system", "temp"):
        return FileSignals(
            file_type=category,
            extension=ext,
            confidence_hint=0.1,
            signals_present=["detected_category"],
        )

    result = FileSignals(
        file_type=category,
        extension=ext,
        signals_present=["detected_category"],
    )

    try:
        if ext == ".pdf":
            _extract_pdf(file_path, result)
        elif ext in {".docx", ".doc"}:
            _extract_docx(file_path, result)
        elif ext in {".pptx", ".ppt"}:
            _extract_pptx(file_path, result)
        elif ext in _TEXT_EXTENSIONS:
            _extract_text(file_path, result)
        elif ext in _CODE_EXTENSIONS:
            _extract_code(file_path, result)
        elif ext in _SPREADSHEET_EXTENSIONS or ext in {".csv",}:
            _extract_spreadsheet(file_path, result)
        elif ext in _IMAGE_EXTENSIONS:
            _extract_image(file_path, result)
        elif ext in _AUDIO_EXTENSIONS:
            _extract_audio(file_path, result)
        elif ext in _VIDEO_EXTENSIONS:
            _extract_video(file_path, result)
        elif ext in _ARCHIVE_EXTENSIONS:
            _extract_archive(file_path, result)
        elif ext in {".rtf", ".odt"}:
            _extract_text_fallback(file_path, result)
        else:
            _extract_text_fallback(file_path, result)
    except Exception as exc:
        logger.debug("Content extraction failed for %s: %s", file_path, exc)

    return result


def _extract_pdf(file_path: Path, result: FileSignals) -> None:
    from pypdf import PdfReader
    reader = PdfReader(str(file_path))
    text_parts = []
    for page in reader.pages:
        try:
            page_text = page.extract_text(extraction_mode="layout") or ""
            if not page_text.strip():
                page_text = page.extract_text() or ""
            if not page_text.strip():
                page_text = _extract_pdf_annotations(page)
        except Exception:
            page_text = ""
        text_parts.append(page_text)
    result.text = "\n".join(text_parts)
    if reader.metadata:
        info = reader.metadata
        if info.title:
            result.metadata["title"] = info.title
        if info.author:
            result.metadata["author"] = info.author
        if info.subject:
            result.metadata["subject"] = info.subject
        if info.creator:
            result.metadata["creator"] = info.creator
        if info.creation_date:
            result.metadata["created_date"] = info.creation_date.strftime("%Y-%m-%d")
    if result.text.strip():
        result.signals_present.append("text")
        result.confidence_hint = max(result.confidence_hint, 0.7)
    if result.metadata:
        result.signals_present.append("metadata")


def _extract_pdf_annotations(page) -> str:
    texts = []
    try:
        for annot in page.get("/Annots", []):
            obj = annot.get_object() if hasattr(annot, "get_object") else annot
            if isinstance(obj, dict):
                for key in ("/V", "/Contents", "/T"):
                    val = obj.get(key)
                    if isinstance(val, str) and val.strip():
                        texts.append(val.strip())
    except Exception:
        pass
    return "\n".join(texts)


def _extract_docx(file_path: Path, result: FileSignals) -> None:
    from docx import Document
    doc = Document(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs]
    result.text = "\n".join(paragraphs)
    try:
        props = doc.core_properties
        if props.author:
            result.metadata["author"] = props.author
        if props.title:
            result.metadata["title"] = props.title
        if props.subject:
            result.metadata["subject"] = props.subject
        if props.created:
            result.metadata["created_date"] = props.created.strftime("%Y-%m-%d")
    except Exception:
        pass
    if result.text.strip():
        result.signals_present.append("text")
        result.confidence_hint = max(result.confidence_hint, 0.8)
    if result.metadata:
        result.signals_present.append("metadata")


def _extract_pptx(file_path: Path, result: FileSignals) -> None:
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        result.text = "\n".join(texts)
        if result.text.strip():
            result.signals_present.append("text")
            result.confidence_hint = max(result.confidence_hint, 0.7)
    except Exception:
        _extract_text_fallback(file_path, result)


def _extract_text(file_path: Path, result: FileSignals) -> None:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        result.text = f.read()
    if result.text.strip():
        result.signals_present.append("text")
        result.confidence_hint = max(result.confidence_hint, 0.9)


def _extract_code(file_path: Path, result: FileSignals) -> None:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    lines = content.splitlines()
    for line in lines[:200]:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("#!") or stripped.startswith("//!") or stripped.startswith("--!"):
            result.metadata.setdefault("shebang", stripped)
        if re.match(r"^(import |from |using |require|include|#include)", stripped):
            result.metadata.setdefault("imports", []).append(stripped)
    result.text = "\n".join(lines[:500])
    if result.text.strip():
        result.signals_present.append("text")
        result.confidence_hint = max(result.confidence_hint, 0.6)


def _extract_spreadsheet(file_path: Path, result: FileSignals) -> None:
    ext = file_path.suffix.lower()
    if ext == ".csv":
        _extract_csv(file_path, result)
    elif ext in _SPREADSHEET_EXTENSIONS:
        _extract_xlsx(file_path, result)


def _extract_xlsx(file_path: Path, result: FileSignals) -> None:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheet_names = wb.sheetnames
        result.metadata["sheet_names"] = sheet_names
        texts = []
        for sheet_name in sheet_names[:3]:
            ws = wb[sheet_name]
            texts.append(f"Sheet: {sheet_name}")
            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if row_idx >= 20:
                    break
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    texts.append(row_text)
                if row_idx == 0:
                    result.metadata["headers"] = [str(c) for c in row if c is not None]
        wb.close()
        result.text = "\n".join(texts)
        if result.text.strip():
            result.signals_present.append("text")
            result.signals_present.append("spreadsheet_data")
            result.confidence_hint = max(result.confidence_hint, 0.6)
    except Exception as exc:
        logger.debug("XLSX extraction failed for %s: %s", file_path, exc)


def _extract_csv(file_path: Path, result: FileSignals) -> None:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            sample = f.read(65536)
        result.text = sample
        if result.text.strip():
            result.signals_present.append("text")
            result.confidence_hint = max(result.confidence_hint, 0.8)
    except Exception as exc:
        logger.debug("CSV extraction failed for %s: %s", file_path, exc)


def _extract_image(file_path: Path, result: FileSignals) -> None:
    try:
        from PIL import Image
        from PIL.ExifTags import TAGS
        img = Image.open(file_path)
        result.metadata["image_format"] = img.format or ""
        result.metadata["width"] = str(img.width)
        result.metadata["height"] = str(img.height)
        result.metadata["mode"] = img.mode
        exif = img.getexif()
        if exif:
            for tag_id, value in exif.items():
                tag_name = TAGS.get(tag_id, str(tag_id))
                if isinstance(value, bytes):
                    try:
                        value = value.decode("utf-8", errors="replace")
                    except Exception:
                        continue
                if tag_name in ("DateTimeOriginal", "DateTime", "CreateDate"):
                    result.metadata["date_taken"] = str(value)
                elif tag_name in ("Make",):
                    result.metadata["camera_make"] = str(value)
                elif tag_name in ("Model",):
                    result.metadata["camera_model"] = str(value)
                elif tag_name in ("GPSInfo",):
                    result.metadata["gps_present"] = "true"
                elif tag_name in ("Software",):
                    result.metadata["software"] = str(value)
                elif tag_name in ("Artist", "Copyright"):
                    result.metadata.setdefault(tag_name.lower(), str(value))
        img.close()
        if result.metadata:
            result.signals_present.append("metadata")
            if "date_taken" in result.metadata or "camera_model" in result.metadata:
                result.confidence_hint = max(result.confidence_hint, 0.5)
    except Exception as exc:
        logger.debug("Image extraction failed for %s: %s", file_path, exc)


def _extract_audio(file_path: Path, result: FileSignals) -> None:
    try:
        from mutagen import File as MutagenFile
        audio = MutagenFile(file_path, easy=True)
        if audio is None:
            return
        if audio.info:
            result.metadata["duration_seconds"] = str(round(audio.info.length, 1))
            if hasattr(audio.info, "bitrate"):
                result.metadata["bitrate"] = str(audio.info.bitrate)
            if hasattr(audio.info, "sample_rate"):
                result.metadata["sample_rate"] = str(audio.info.sample_rate)
        if audio.tags:
            for key in ("title", "artist", "album", "date", "genre", "tracknumber"):
                val = audio.tags.get(key)
                if val:
                    result.metadata[key] = val[0] if isinstance(val, list) else str(val)
        if result.metadata:
            result.signals_present.append("metadata")
            if "artist" in result.metadata or "title" in result.metadata:
                result.confidence_hint = max(result.confidence_hint, 0.6)
    except Exception as exc:
        logger.debug("Audio extraction failed for %s: %s", file_path, exc)


def _extract_video(file_path: Path, result: FileSignals) -> None:
    try:
        from pymediainfo import MediaInfo
        media = MediaInfo.parse(file_path)
        for track in media.tracks:
            if track.track_type == "General":
                if track.duration:
                    result.metadata["duration_ms"] = track.duration
                if track.file_size:
                    result.metadata["file_size_bytes"] = str(track.file_size)
            elif track.track_type == "Video":
                if track.format:
                    result.metadata["video_format"] = track.format
                if track.width and track.height:
                    result.metadata["resolution"] = f"{track.width}x{track.height}"
                if track.frame_rate:
                    result.metadata["frame_rate"] = track.frame_rate
            elif track.track_type == "Audio":
                if track.format:
                    result.metadata["audio_format"] = track.format
                if track.channel_s:
                    result.metadata["audio_channels"] = str(track.channel_s)
        if result.metadata:
            result.signals_present.append("metadata")
            result.confidence_hint = max(result.confidence_hint, 0.4)
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("Video extraction failed for %s: %s", file_path, exc)


def _extract_archive(file_path: Path, result: FileSignals) -> None:
    ext = file_path.suffix.lower()
    try:
        names = []
        if ext == ".zip":
            with zipfile.ZipFile(file_path, "r") as zf:
                names = zf.namelist()
        elif ext in {".tar", ".gz", ".bz2", ".xz", ".zst"}:
            mode = "r:*"
            with tarfile.open(file_path, mode) as tf:
                names = [m.name for m in tf.getmembers()]
        if names:
            result.metadata["archive_contents"] = names[:50]
            result.metadata["archive_count"] = str(len(names))
            result.signals_present.append("archive_listing")
            result.confidence_hint = max(result.confidence_hint, 0.3)
    except Exception as exc:
        logger.debug("Archive extraction failed for %s: %s", file_path, exc)


def _extract_text_fallback(file_path: Path, result: FileSignals) -> None:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            result.text = f.read(65536)
        if result.text.strip():
            result.signals_present.append("text")
            result.confidence_hint = max(result.confidence_hint, 0.4)
    except Exception:
        pass
