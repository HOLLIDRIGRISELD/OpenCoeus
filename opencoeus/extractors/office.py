from __future__ import annotations

import logging
import re
from pathlib import Path

from .signals import FileSignals

logger = logging.getLogger(__name__)

_DOC_TYPE_PATTERNS: list[tuple[str, re.Pattern]] = [
    ("Invoice", re.compile(r"(invoice|inv[-#\s]?\d|amount\s+d(ue|ate)|billing|payment\s+terms)", re.IGNORECASE)),
    ("Meeting-Notes", re.compile(r"(meeting\s*(notes|minutes|log)|attendees?:|agenda|action\s+items)", re.IGNORECASE)),
    ("Specification", re.compile(r"(specification|spec[:.\s]|technical\s*(spec|requirement)|architecture\s*(overview|diagram))", re.IGNORECASE)),
    ("Report", re.compile(r"(report|summary|analysis|findings|conclusion)", re.IGNORECASE)),
    ("Budget", re.compile(r"(budget|financial|revenue|expenses?|forecast|quarterly)", re.IGNORECASE)),
    ("Contract", re.compile(r"(contract|agreement|terms?\s*(and|of)\s*(service|conditions|use)|legal)", re.IGNORECASE)),
    ("Presentation", re.compile(r"(presentation|slides?|deck)", re.IGNORECASE)),
    ("Readme", re.compile(r"(readme|getting\s+started)", re.IGNORECASE)),
    ("Backup", re.compile(r"(backup|archive|restore|snapshot)", re.IGNORECASE)),
    ("Spreadsheet", re.compile(r"(spreadsheet|worksheet|tab|csv|tabular)", re.IGNORECASE)),
    ("Resume", re.compile(r"(resume|cv|curriculum\s+vitae)", re.IGNORECASE)),
    ("Letter", re.compile(r"(dear\s+\w+|sincerely|yours\s+faithfully)", re.IGNORECASE)),
    ("Memo", re.compile(r"(memo|memorandum|to:?|from:?|re:?|subject:?)", re.IGNORECASE)),
    ("Article", re.compile(r"(article|abstract|introduction|methodology|references)", re.IGNORECASE)),
    ("Proposal", re.compile(r"(proposal|statement\s+of\s+work|scope\s+of\s+work|deliverables)", re.IGNORECASE)),
    ("Manual", re.compile(r"(manual|guide|tutorial|user\s+guide|reference)", re.IGNORECASE)),
    ("Certificate", re.compile(r"(certificate|certification|completion|certified)", re.IGNORECASE)),
    ("Form", re.compile(r"(form|application|registration|questionnaire)", re.IGNORECASE)),
    ("Policy", re.compile(r"(policy|procedure|compliance|regulation)", re.IGNORECASE)),
    ("Checklist", re.compile(r"(checklist|check\s*list|todo|tasks)", re.IGNORECASE)),
    ("Thesis", re.compile(r"(thesis|dissertation|doctoral|master['\u2019]s\s+thesis)", re.IGNORECASE)),
    ("Agenda", re.compile(r"(agenda|schedule|timeline|itinerary)", re.IGNORECASE)),
    ("Newsletter", re.compile(r"(newsletter|news\s*letter|issue\s+\d)", re.IGNORECASE)),
    ("Order", re.compile(r"(order|purchase\s+order|po\s*[-\s]?\d+)", re.IGNORECASE)),
    ("Quote", re.compile(r"(quote|quotation|estimate|price\s+list)", re.IGNORECASE)),
    ("Timesheet", re.compile(r"(timesheet|time\s*sheet|time\s+tracking|hours)", re.IGNORECASE)),
]

_NOISE_PATTERNS = re.compile(
    r"^(page\s+\d+|第.+页|header|footer|confidential|draft|do not distribute)$",
    re.IGNORECASE,
)

_TITLE_INDICATORS = re.compile(
    r"(title|subject|report|invoice|specification|summary|meeting|notes|agenda|minutes)",
    re.IGNORECASE,
)

_HEADER_PREFIXES = re.compile(
    r"^(#+\s+|[-=]+\s*$|abstract|introduction|overview|background)",
    re.IGNORECASE,
)

_CODE_PATTERNS = re.compile(r"[{}();\[\]\\]")
_URL_PATTERN = re.compile(r"https?://|www\.|ftp://")
_PATH_PATTERN = re.compile(r"^[/\\]|[A-Za-z]:[/\\]")
_PROPER_SENTENCE = re.compile(r"^[A-Z][a-z]")


def _extract_docx(file_path: Path, result: FileSignals) -> None:
    from docx import Document
    doc = Document(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs]
    result.text = "\n".join(paragraphs)
    try:
        props = doc.core_properties
        if props.author:
            result.metadata["author"] = props.author
        if props.title:
            result.metadata["title"] = props.title
        if props.subject:
            result.metadata["subject"] = props.subject
        if props.created:
            result.metadata["created_date"] = props.created.strftime("%Y-%m-%d")
    except Exception:
        pass
    if result.text.strip():
        result.signals_present.append("text")
        result.confidence_hint = max(result.confidence_hint, 0.8)
    if result.metadata:
        result.signals_present.append("metadata")


def _extract_pptx(file_path: Path, result: FileSignals) -> None:
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        result.text = "\n".join(texts)
        if result.text.strip():
            result.signals_present.append("text")
            result.confidence_hint = max(result.confidence_hint, 0.7)
    except Exception:
        from .base import _extract_text_fallback
        _extract_text_fallback(file_path, result)


def _extract_spreadsheet(file_path: Path, result: FileSignals) -> None:
    ext = file_path.suffix.lower()
    if ext == ".csv":
        _extract_csv(file_path, result)
    elif ext in {".xlsx", ".xls", ".xlsm", ".xlsb", ".ods"}:
        _extract_xlsx(file_path, result)


def _extract_xlsx(file_path: Path, result: FileSignals) -> None:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheet_names = wb.sheetnames
        result.metadata["sheet_names"] = sheet_names
        texts = []
        for sheet_name in sheet_names[:3]:
            ws = wb[sheet_name]
            texts.append(f"Sheet: {sheet_name}")
            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if row_idx >= 20:
                    break
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    texts.append(row_text)
                if row_idx == 0:
                    result.metadata["headers"] = [str(c) for c in row if c is not None]
        wb.close()
        result.text = "\n".join(texts)
        if result.text.strip():
            result.signals_present.append("text")
            result.signals_present.append("spreadsheet_data")
            result.confidence_hint = max(result.confidence_hint, 0.6)
    except Exception as exc:
        logger.debug("XLSX extraction failed for %s: %s", file_path, exc)


def _extract_csv(file_path: Path, result: FileSignals) -> None:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            sample = f.read(65536)
        result.text = sample
        if result.text.strip():
            result.signals_present.append("text")
            result.confidence_hint = max(result.confidence_hint, 0.8)
    except Exception as exc:
        logger.debug("CSV extraction failed for %s: %s", file_path, exc)


# --- functions moved from opencoeus/documents.py ---


def _extract_pdf_text_annotations(pdf_page) -> str:
    texts: list[str] = []
    try:
        for annot in pdf_page.get("/Annots", []):
            annot_obj = annot.get_object() if hasattr(annot, "get_object") else annot
            if isinstance(annot_obj, dict):
                for key in ("/V", "/Contents", "/T"):
                    value = annot_obj.get(key)
                    if isinstance(value, str) and value.strip():
                        texts.append(value.strip())
    except Exception:
        pass
    return "\n".join(texts)


def extract_text(document_path: Path, maximum_pages: int = 2) -> str:
    logger.debug("Extracting text from %s", document_path)
    try:
        if document_path.suffix.lower() == ".pdf":
            from pypdf import PdfReader
            pdf_reader = PdfReader(str(document_path))
            text_parts: list[str] = []
            for pdf_page in pdf_reader.pages[:maximum_pages] if maximum_pages < 999 else pdf_reader.pages:
                page_text = pdf_page.extract_text(extraction_mode="layout") or ""
                if not page_text.strip():
                    page_text = pdf_page.extract_text() or ""
                if not page_text.strip():
                    page_text = _extract_pdf_text_annotations(pdf_page)
                text_parts.append(page_text)
            return "\n".join(text_parts)
        if document_path.suffix.lower() == ".docx":
            from docx import Document
            word_document = Document(str(document_path))
            paragraphs = word_document.paragraphs[:80] if maximum_pages < 999 else word_document.paragraphs
            return "\n".join(p.text for p in paragraphs)
        if document_path.suffix.lower() in {".txt", ".md"}:
            lines: list[str] = []
            with open(document_path, "r", encoding="utf-8", errors="replace") as text_file:
                if maximum_pages < 999:
                    for line_number, line in enumerate(text_file):
                        if line_number >= 15:
                            break
                        lines.append(line)
                else:
                    lines = text_file.readlines()
            return "".join(lines)
        if document_path.suffix.lower() in {".xlsx", ".xls"}:
            return _extract_xlsx_text(document_path)
        if document_path.suffix.lower() == ".csv":
            return _extract_csv_text(document_path)
    except Exception as exc:
        logger.debug("Text extraction failed for %s: %s", document_path, exc)
        return ""
    return ""


def _extract_xlsx_text(file_path: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames[:3]:
            ws = wb[sheet_name]
            texts.append(f"Sheet: {sheet_name}")
            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if row_idx >= 20:
                    break
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    texts.append(row_text)
        wb.close()
        return "\n".join(texts)
    except Exception:
        return ""


def _extract_csv_text(file_path: Path) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(65536)
    except Exception:
        return ""


def extract_metadata(document_path: Path) -> dict:
    metadata: dict = {}
    try:
        if document_path.suffix.lower() == ".pdf":
            from pypdf import PdfReader
            pdf_reader = PdfReader(str(document_path))
            info = pdf_reader.metadata
            if info:
                metadata["author"] = info.author or ""
                metadata["title"] = info.title or ""
                metadata["subject"] = info.subject or ""
                metadata["creator"] = info.creator or ""
                if info.creation_date:
                    metadata["created_date"] = info.creation_date.strftime("%Y-%m-%d")
                    metadata["created_month"] = info.creation_date.strftime("%m")
        if document_path.suffix.lower() == ".docx":
            from docx import Document
            word_document = Document(str(document_path))
            try:
                props = word_document.core_properties
                if props.author:
                    metadata["author"] = props.author
                if props.title:
                    metadata["title"] = props.title
                if props.subject:
                    metadata["subject"] = props.subject
                if props.created:
                    metadata["created_date"] = props.created.strftime("%Y-%m-%d")
                    metadata["created_month"] = props.created.strftime("%m")
            except Exception:
                pass
    except Exception as exc:
        logger.debug("Metadata extraction failed for %s: %s", document_path, exc)
    return metadata


def _score_title_candidate(line: str) -> float:
    score = 0.0
    stripped = line.strip()
    words = stripped.split()
    word_count = len(words)

    if _NOISE_PATTERNS.match(stripped):
        score -= 3.0
    if _HEADER_PREFIXES.match(stripped):
        score -= 1.0
    if _CODE_PATTERNS.search(stripped):
        score -= 2.0
    if _URL_PATTERN.search(stripped):
        score -= 3.0
    if _PATH_PATTERN.match(stripped):
        score -= 2.0

    if 2 <= word_count <= 8:
        score += 2.0
    elif 1 <= word_count <= 12:
        score += 1.0

    if any(c.isupper() for c in stripped) and any(c.islower() for c in stripped):
        score += 1.5
    elif stripped.isupper():
        score -= 0.5

    if _PROPER_SENTENCE.match(stripped):
        score += 0.5

    alpha_ratio = sum(c.isalpha() for c in stripped) / max(len(stripped), 1)
    if alpha_ratio > 0.7:
        score += 1.0
    elif alpha_ratio > 0.4:
        score += 0.3

    if _TITLE_INDICATORS.search(stripped):
        score += 2.0

    if 20 <= len(stripped) <= 80:
        score += 0.5

    return score


def suggest_title(document_text: str, fallback_title: str) -> str:
    text_lines = [re.sub(r"\s+", " ", line).strip() for line in document_text.splitlines()]
    candidate_titles = [
        line for line in text_lines
        if 5 <= len(line) <= 120 and any(character.isalpha() for character in line)
    ]
    if not candidate_titles:
        safe = re.sub(r'[<>:"/\\|?*]', "-", fallback_title).strip(" .-")
        return safe[:120] or "Untitled document"

    scored = [(_score_title_candidate(line), line) for line in candidate_titles]
    best_score, best_line = max(scored, key=lambda x: x[0])
    if best_score < 1.0:
        safe = re.sub(r'[<>:"/\\|?*]', "-", fallback_title).strip(" .-")
        return safe[:120] or "Untitled document"

    safe = re.sub(r'[<>:"/\\|?*]', "-", best_line).strip(" .-")
    safe = re.sub(r"\s+", " ", safe)
    return safe[:120] or "Untitled document"


def detect_document_type(text: str, metadata: dict | None = None, title_candidate: str = "") -> str:
    combined = text
    if title_candidate:
        combined = title_candidate + "\n" + text
    if metadata:
        meta_text = " ".join(str(v) for v in metadata.values() if v)
        combined = combined + "\n" + meta_text
    for doc_type_name, doc_type_pattern in _DOC_TYPE_PATTERNS:
        if doc_type_pattern.search(combined):
            return doc_type_name
    return "Document"
